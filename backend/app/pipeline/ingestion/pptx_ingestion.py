from pptx import Presentation
from app.pipeline.step import Step
import pandas as pd


class PptxIngestion(Step):
    def __init__(self, file_path):
        self.file_path = file_path

    def execute(self, data={}):
        prs = Presentation(self.file_path)
        slides_data = []
        for slide in prs.slides:
            slide_data = {"texts": [], "tables": []}
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_data["texts"].append(shape.text)

                if shape.has_table:
                    table = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table.append(row_data)

                    df = pd.DataFrame(
                        table[1:], columns=table[0]
                    )  # Use the first row as headers

                    slide_data["tables"].append(df.to_dict(orient="records"))
            slides_data.append(slide_data)
        return slides_data
